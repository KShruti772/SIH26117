import os
import re
import csv
import io
import hashlib
import logging
from dataclasses import dataclass, field
from typing import List, Dict, Any, Optional
from PIL import Image

from backend.rag.detector import FileDetector, DetectionResult
from backend.multimodal.ocr import BaseOCR, LocalPytesseractOCR

logger = logging.getLogger("aegis.rag.extractors")

class InsufficientTextError(Exception):
    """Raised when document extraction yields insufficient or no text."""
    pass

@dataclass
class NormalizedPage:
    """Represents a single logical page, slide, sheet, or code section within a document."""
    page_number: int
    section_title: str
    text: str
    table_data: Optional[List[List[str]]] = None
    metadata: Dict[str, Any] = field(default_factory=dict)

    def to_dict(self) -> Dict[str, Any]:
        return {
            "page_number": self.page_number,
            "section": self.section_title,
            "text": self.text,
            "table_data": self.table_data,
            "metadata": self.metadata
        }

@dataclass
class NormalizedDocument:
    """Universal internal normalized representation for all ingested documents."""
    doc_id: str
    filename: str
    category: str
    file_type: str
    mime_type: str
    file_size: int
    content_hash: str
    pages: List[NormalizedPage]
    extraction_method: str = "native"
    metadata: Dict[str, Any] = field(default_factory=dict)

    def get_total_text(self) -> str:
        return "\n\n".join(p.text for p in self.pages if p.text.strip())

    def to_dict(self) -> Dict[str, Any]:
        return {
            "id": self.doc_id,
            "filename": self.filename,
            "category": self.category,
            "file_type": self.file_type,
            "mime_type": self.mime_type,
            "file_size": self.file_size,
            "content_hash": self.content_hash,
            "chunk_count": len(self.pages),
            "extraction_method": self.extraction_method,
            "metadata": self.metadata,
            "pages": [p.to_dict() for p in self.pages]
        }

class BaseExtractor:
    """Base interface for all file category extractors."""
    def extract(
        self,
        file_path: str,
        filename: str,
        detection: DetectionResult,
        ocr_service: Optional[BaseOCR] = None
    ) -> NormalizedDocument:
        raise NotImplementedError

class PDFExtractor(BaseExtractor):
    """Extracts digital and scanned PDF pages with automatic OCR failover."""
    def extract(
        self,
        file_path: str,
        filename: str,
        detection: DetectionResult,
        ocr_service: Optional[BaseOCR] = None
    ) -> NormalizedDocument:
        import fitz
        pages: List[NormalizedPage] = []
        total_text = ""
        doc_metadata: Dict[str, Any] = {}

        try:
            doc = fitz.open(file_path)
            doc_metadata = {
                "page_count": len(doc),
                "title": doc.metadata.get("title", ""),
                "author": doc.metadata.get("author", "")
            }

            for idx, page in enumerate(doc):
                text = page.get_text().strip()
                total_text += text
                pages.append(NormalizedPage(
                    page_number=idx + 1,
                    section_title=f"Page {idx + 1}",
                    text=text,
                    metadata={"page_number": idx + 1}
                ))
            doc.close()
        except Exception as e:
            logger.warning(f"PyMuPDF failed on {filename}: {e}. Trying pypdf fallback.")
            try:
                import pypdf
                reader = pypdf.PdfReader(file_path)
                pages = []
                total_text = ""
                for idx, page in enumerate(reader.pages):
                    t = (page.extract_text() or "").strip()
                    total_text += t
                    pages.append(NormalizedPage(
                        page_number=idx + 1,
                        section_title=f"Page {idx + 1}",
                        text=t,
                        metadata={"page_number": idx + 1}
                    ))
            except Exception as pe:
                raise ValueError(f"Failed to parse PDF document contents: {pe}")

        # Check for scanned/image-only PDF requiring local OCR
        extraction_method = "native_pdf"
        if len(total_text.strip()) < 15:
            logger.info(f"PDF '{filename}' has insufficient digital text ({len(total_text)} chars). Invoking local OCR.")
            ocr = ocr_service or LocalPytesseractOCR()
            if ocr.is_available():
                ocr_res = ocr.ocr_pdf(file_path)
                pages = [
                    NormalizedPage(
                        page_number=p.get("page_number", idx + 1),
                        section_title=f"Page {p.get('page_number', idx + 1)} (OCR)",
                        text=p.get("text", "").strip(),
                        metadata={"ocr_processed": True, "page_number": p.get("page_number", idx + 1)}
                    )
                    for idx, p in enumerate(ocr_res.get("pages", []))
                    if p.get("text", "").strip()
                ]
                extraction_method = "ocr_scanned_pdf"
            else:
                if not pages or not any(p.text.strip() for p in pages):
                    raise InsufficientTextError("Scanned PDF contains no digital text and local OCR engine is unavailable.")

        with open(file_path, "rb") as f:
            content = f.read()
            chash = hashlib.sha256(content).hexdigest()

        return NormalizedDocument(
            doc_id=chash,
            filename=filename,
            category="document",
            file_type="pdf",
            mime_type="application/pdf",
            file_size=len(content),
            content_hash=chash,
            pages=pages,
            extraction_method=extraction_method,
            metadata=doc_metadata
        )

class DocxExtractor(BaseExtractor):
    """Extracts Word (.docx) paragraphs, headings, lists, tables, and section hierarchies."""
    def extract(
        self,
        file_path: str,
        filename: str,
        detection: DetectionResult,
        ocr_service: Optional[BaseOCR] = None
    ) -> NormalizedDocument:
        import docx
        doc = docx.Document(file_path)

        pages: List[NormalizedPage] = []
        current_text: List[str] = []
        current_page = 1
        current_section = "Document Body"
        current_len = 0
        page_size = 2000

        for p in doc.paragraphs:
            text = p.text.strip()
            if not text:
                continue
            if p.style and "Heading" in p.style.name:
                current_section = text

            if current_len + len(text) > page_size and current_text:
                pages.append(NormalizedPage(
                    page_number=current_page,
                    section_title=current_section,
                    text="\n\n".join(current_text),
                    metadata={"section": current_section, "page_number": current_page}
                ))
                current_page += 1
                current_text = [text]
                current_len = len(text)
            else:
                current_text.append(text)
                current_len += len(text)

        # Extract structured tables
        for tbl_idx, table in enumerate(doc.tables):
            table_lines: List[str] = []
            table_grid: List[List[str]] = []
            for row in table.rows:
                cells = [c.text.strip() for c in row.cells]
                if any(cells):
                    table_lines.append(" | ".join(cells))
                    table_grid.append(cells)
            if table_lines:
                tbl_text = f"### Table {tbl_idx + 1}\n" + "\n".join(table_lines)
                current_text.append(tbl_text)
                current_len += len(tbl_text)

        if current_text:
            pages.append(NormalizedPage(
                page_number=current_page,
                section_title=current_section,
                text="\n\n".join(current_text),
                metadata={"section": current_section, "page_number": current_page}
            ))

        if not pages or not any(p.text.strip() for p in pages):
            raise ValueError("DOCX document contains no extractable text.")

        with open(file_path, "rb") as f:
            content = f.read()
            chash = hashlib.sha256(content).hexdigest()

        return NormalizedDocument(
            doc_id=chash,
            filename=filename,
            category="document",
            file_type="docx",
            mime_type=detection.mime_type,
            file_size=len(content),
            content_hash=chash,
            pages=pages,
            extraction_method="docx_parser",
            metadata={"paragraph_count": len(doc.paragraphs), "table_count": len(doc.tables)}
        )

class SpreadsheetExtractor(BaseExtractor):
    """Extracts Excel (XLSX, XLS) and Tabular (CSV, TSV) workbooks preserving sheet and table structure."""
    def extract(
        self,
        file_path: str,
        filename: str,
        detection: DetectionResult,
        ocr_service: Optional[BaseOCR] = None
    ) -> NormalizedDocument:
        pages: List[NormalizedPage] = []
        ext = detection.extension.lower()

        if ext in [".xlsx", ".xlsm"]:
            import openpyxl
            wb = openpyxl.load_workbook(file_path, data_only=True, read_only=True)
            page_num = 1
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                rows = []
                for row in sheet.iter_rows(values_only=True):
                    row_vals = [str(cell) if cell is not None else "" for cell in row]
                    if any(v.strip() for v in row_vals):
                        rows.append(row_vals)

                if not rows:
                    continue

                header = " | ".join(rows[0])
                batch_size = 35
                for i in range(1, max(2, len(rows)), batch_size):
                    batch = rows[i:i + batch_size]
                    lines = [f"### Sheet: {sheet_name}", header, "---"]
                    for r in batch:
                        lines.append(" | ".join(r))

                    section_name = f"Sheet: {sheet_name}"
                    pages.append(NormalizedPage(
                        page_number=page_num,
                        section_title=section_name,
                        text="\n".join(lines),
                        table_data=[rows[0]] + batch,
                        metadata={"sheet_name": sheet_name, "row_start": i, "row_end": min(i+batch_size-1, len(rows)-1)}
                    ))
                    page_num += 1
            wb.close()

        else:
            # CSV or TSV
            delimiter = "\t" if ext in [".tsv", ".tab"] else ","
            rows = []
            try:
                with open(file_path, "r", encoding="utf-8") as f:
                    reader = csv.reader(f, delimiter=delimiter)
                    rows = [r for r in reader if any(c.strip() for c in r)]
            except UnicodeDecodeError:
                with open(file_path, "r", encoding="latin-1") as f:
                    reader = csv.reader(f, delimiter=delimiter)
                    rows = [r for r in reader if any(c.strip() for c in r)]

            if not rows:
                raise ValueError("Tabular file is empty.")

            header = " | ".join(rows[0])
            batch_size = 30
            page_num = 1
            for i in range(1, max(2, len(rows)), batch_size):
                batch = rows[i:i + batch_size]
                lines = [header, "---"]
                for r in batch:
                    lines.append(" | ".join(r))

                row_end = min(i + batch_size - 1, len(rows) - 1)
                section_name = f"Rows {i}-{row_end}" if len(rows) > 1 else "Header"
                pages.append(NormalizedPage(
                    page_number=page_num,
                    section_title=section_name,
                    text="\n".join(lines),
                    table_data=[rows[0]] + batch,
                    metadata={"row_start": i, "row_end": row_end}
                ))
                page_num += 1

        if not pages:
            raise ValueError("Spreadsheet document contains no readable tabular records.")

        with open(file_path, "rb") as f:
            content = f.read()
            chash = hashlib.sha256(content).hexdigest()

        return NormalizedDocument(
            doc_id=chash,
            filename=filename,
            category="spreadsheet",
            file_type=detection.file_type,
            mime_type=detection.mime_type,
            file_size=len(content),
            content_hash=chash,
            pages=pages,
            extraction_method="tabular",
            metadata={"format": detection.file_type}
        )

class PresentationExtractor(BaseExtractor):
    """Extracts PowerPoint (.pptx, .odp) slides, titles, shapes, tables, and notes."""
    def extract(
        self,
        file_path: str,
        filename: str,
        detection: DetectionResult,
        ocr_service: Optional[BaseOCR] = None
    ) -> NormalizedDocument:
        import pptx
        prs = pptx.Presentation(file_path)
        pages: List[NormalizedPage] = []

        for idx, slide in enumerate(prs.slides):
            slide_num = idx + 1
            slide_title = f"Slide {slide_num}"
            slide_texts: List[str] = []

            # Extract title if present
            if slide.shapes.title and slide.shapes.title.text.strip():
                slide_title = f"Slide {slide_num}: {slide.shapes.title.text.strip()}"
                slide_texts.append(f"## {slide.shapes.title.text.strip()}")

            # Extract body text from shapes
            for shape in slide.shapes:
                if shape != slide.shapes.title and shape.has_text_frame:
                    t = shape.text_frame.text.strip()
                    if t:
                        slide_texts.append(t)
                elif shape.has_table:
                    table_rows = []
                    for row in shape.table.rows:
                        cells = [c.text.strip() for c in row.cells]
                        if any(cells):
                            table_rows.append(" | ".join(cells))
                    if table_rows:
                        slide_texts.append("\n".join(table_rows))

            # Extract speaker notes
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    slide_texts.append(f"**Speaker Notes:** {notes}")

            combined = "\n\n".join(slide_texts).strip()
            if not combined:
                combined = f"Slide {slide_num} [Visual layout / Empty text]"

            pages.append(NormalizedPage(
                page_number=slide_num,
                section_title=f"Slide {slide_num}",
                text=combined,
                metadata={"slide_number": slide_num, "slide_title": slide_title}
            ))

        if not pages:
            raise ValueError("Presentation contains no slides.")

        with open(file_path, "rb") as f:
            content = f.read()
            chash = hashlib.sha256(content).hexdigest()

        return NormalizedDocument(
            doc_id=chash,
            filename=filename,
            category="presentation",
            file_type="pptx",
            mime_type=detection.mime_type,
            file_size=len(content),
            content_hash=chash,
            pages=pages,
            extraction_method="slide_parser",
            metadata={"slide_count": len(prs.slides)}
        )

class ImageMultimodalExtractor(BaseExtractor):
    """Extracts OCR text, visual metadata, and prepares image representations."""
    def extract(
        self,
        file_path: str,
        filename: str,
        detection: DetectionResult,
        ocr_service: Optional[BaseOCR] = None
    ) -> NormalizedDocument:
        with Image.open(file_path) as img:
            width, height = img.size
            img_format = img.format or detection.file_type.upper()
            mode = img.mode

        ocr = ocr_service or LocalPytesseractOCR()
        ocr_text = ""
        if ocr.is_available():
            try:
                ocr_text = ocr.ocr_image(file_path).strip()
            except Exception as e:
                logger.warning(f"OCR image extraction warning on {filename}: {e}")

        visual_summary = (
            f"**Image / Visual Artifact:** {filename}\n"
            f"**Dimensions:** {width}x{height} pixels | **Format:** {img_format} | **Color Mode:** {mode}\n"
        )
        if ocr_text:
            visual_summary += f"\n**Extracted Text / Labels (OCR):**\n{ocr_text}"
        else:
            visual_summary += "\n**Extracted Text:** No text detected on image."

        pages = [
            NormalizedPage(
                page_number=1,
                section_title=f"Image: {filename}",
                text=visual_summary,
                metadata={
                    "width": width,
                    "height": height,
                    "format": img_format,
                    "has_ocr_text": bool(ocr_text)
                }
            )
        ]

        with open(file_path, "rb") as f:
            content = f.read()
            chash = hashlib.sha256(content).hexdigest()

        return NormalizedDocument(
            doc_id=chash,
            filename=filename,
            category="image",
            file_type=detection.file_type,
            mime_type=detection.mime_type,
            file_size=len(content),
            content_hash=chash,
            pages=pages,
            extraction_method="vision_ocr",
            metadata={"width": width, "height": height, "format": img_format}
        )

class CodeExtractor(BaseExtractor):
    """Extracts programming language files (Python, JS, TS, Java, C/C++, SQL, Shell, YAML, JSON, XML)."""
    def extract(
        self,
        file_path: str,
        filename: str,
        detection: DetectionResult,
        ocr_service: Optional[BaseOCR] = None
    ) -> NormalizedDocument:
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                lines = f.readlines()
        except UnicodeDecodeError:
            with open(file_path, "r", encoding="latin-1") as f:
                lines = f.readlines()

        if not lines:
            raise ValueError("Code file is empty.")

        pages: List[NormalizedPage] = []
        batch_size = 50
        page_num = 1

        for i in range(0, len(lines), batch_size):
            batch_lines = lines[i:i + batch_size]
            line_start = i + 1
            line_end = i + len(batch_lines)
            section_title = f"Lines {line_start}-{line_end}"

            numbered_code = "".join(f"{line_start + idx:4d} | {l}" for idx, l in enumerate(batch_lines))
            code_block = (
                f"```{detection.file_type}\n"
                f"# File: {filename} ({section_title})\n"
                f"{numbered_code}\n"
                f"```"
            )

            pages.append(NormalizedPage(
                page_number=page_num,
                section_title=section_title,
                text=code_block,
                metadata={"line_start": line_start, "line_end": line_end, "language": detection.file_type}
            ))
            page_num += 1

        with open(file_path, "rb") as f:
            content = f.read()
            chash = hashlib.sha256(content).hexdigest()

        return NormalizedDocument(
            doc_id=chash,
            filename=filename,
            category="code",
            file_type=detection.file_type,
            mime_type=detection.mime_type,
            file_size=len(content),
            content_hash=chash,
            pages=pages,
            extraction_method="code_parser",
            metadata={"language": detection.file_type, "total_lines": len(lines)}
        )

class TextExtractor(BaseExtractor):
    """Extracts plain text, markdown, RTF, and log files."""
    def extract(
        self,
        file_path: str,
        filename: str,
        detection: DetectionResult,
        ocr_service: Optional[BaseOCR] = None
    ) -> NormalizedDocument:
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                content = f.read()
        except UnicodeDecodeError:
            with open(file_path, "r", encoding="latin-1") as f:
                content = f.read()

        if not content.strip():
            raise InsufficientTextError("Text document is empty.")

        pages: List[NormalizedPage] = []
        ext = detection.extension.lower()

        if ext in [".md", ".markdown"]:
            sections = re.split(r'\n(?=#{1,3}\s+)', content)
            page_num = 1
            for sec in sections:
                sec_clean = sec.strip()
                if not sec_clean:
                    continue
                first_line = sec_clean.split("\n", 1)[0].lstrip("#").strip()
                pages.append(NormalizedPage(
                    page_number=page_num,
                    section_title=first_line or f"Section {page_num}",
                    text=sec_clean,
                    metadata={"section": first_line, "page_number": page_num}
                ))
                page_num += 1
        else:
            paragraphs = content.split("\n\n")
            current_page_text: List[str] = []
            current_len = 0
            page_size = 2500
            page_num = 1

            for p in paragraphs:
                p_str = p.strip()
                if not p_str:
                    continue
                if current_len + len(p_str) > page_size and current_page_text:
                    pages.append(NormalizedPage(
                        page_number=page_num,
                        section_title=f"Section {page_num}",
                        text="\n\n".join(current_page_text),
                        metadata={"page_number": page_num}
                    ))
                    page_num += 1
                    current_page_text = [p_str]
                    current_len = len(p_str)
                else:
                    current_page_text.append(p_str)
                    current_len += len(p_str)

            if current_page_text:
                pages.append(NormalizedPage(
                    page_number=page_num,
                    section_title=f"Section {page_num}",
                    text="\n\n".join(current_page_text),
                    metadata={"page_number": page_num}
                ))

        with open(file_path, "rb") as f:
            raw_bytes = f.read()
            chash = hashlib.sha256(raw_bytes).hexdigest()

        return NormalizedDocument(
            doc_id=chash,
            filename=filename,
            category="text",
            file_type=detection.file_type,
            mime_type=detection.mime_type,
            file_size=len(raw_bytes),
            content_hash=chash,
            pages=pages if pages else [NormalizedPage(page_number=1, section_title="General", text=content.strip())],
            extraction_method="native",
            metadata={"format": detection.file_type}
        )

class UniversalExtractorRegistry:
    """Central registry resolving the optimal category extractor."""
    _extractors = {
        "pdf": PDFExtractor(),
        "docx": DocxExtractor(),
        "spreadsheet": SpreadsheetExtractor(),
        "presentation": PresentationExtractor(),
        "image": ImageMultimodalExtractor(),
        "code": CodeExtractor(),
        "text": TextExtractor()
    }

    @classmethod
    def get_extractor(cls, detection: DetectionResult) -> BaseExtractor:
        if detection.file_type == "pdf":
            return cls._extractors["pdf"]
        elif detection.file_type == "docx":
            return cls._extractors["docx"]
        elif detection.category == "spreadsheet":
            return cls._extractors["spreadsheet"]
        elif detection.category == "presentation":
            return cls._extractors["presentation"]
        elif detection.category == "image":
            return cls._extractors["image"]
        elif detection.category == "code":
            return cls._extractors["code"]
        else:
            return cls._extractors["text"]

    @classmethod
    def extract_document(
        cls,
        file_path: str,
        filename: Optional[str] = None,
        detection: Optional[DetectionResult] = None,
        ocr_service: Optional[BaseOCR] = None
    ) -> NormalizedDocument:
        target_name = filename or os.path.basename(file_path)
        det = detection or FileDetector.detect_from_path(file_path, filename_override=target_name)

        if not det.is_safe or not det.is_valid:
            raise ValueError(det.error_reason or f"Unsupported or dangerous file type for '{target_name}'.")

        extractor = cls.get_extractor(det)
        return extractor.extract(file_path, target_name, det, ocr_service=ocr_service)
