import os
import io
import json
import zipfile
import tempfile
import unittest
import hashlib
from unittest.mock import MagicMock, patch

from PIL import Image, ImageDraw
import docx
import openpyxl
import pptx

from backend.rag.detector import FileDetector, DetectionResult
from backend.rag.extractors import (
    UniversalExtractorRegistry,
    NormalizedDocument,
    NormalizedPage,
    PDFExtractor,
    DocxExtractor,
    SpreadsheetExtractor,
    PresentationExtractor,
    ImageMultimodalExtractor,
    CodeExtractor,
    TextExtractor
)
from backend.rag.pipeline import AegisRagService
from backend.rag.embeddings import BaseEmbeddingModel
from backend.multimodal.ocr import BaseOCR
from backend.security.database import init_db

class MockEmbeddingModel(BaseEmbeddingModel):
    def __init__(self, dimension: int = 64):
        self.dimension = dimension
        self.model_name = "test-mock-embeddings"
        self.is_mock = True

    def embed_text(self, text: str):
        import math
        import re
        vec = [0.0] * self.dimension
        words = re.findall(r'\w+', text.lower())
        if not words:
            return vec
        for w in words:
            h = int(hashlib.md5(w.encode("utf-8")).hexdigest(), 16) % self.dimension
            vec[h] += 1.0
        norm = math.sqrt(sum(x * x for x in vec))
        if norm > 0:
            vec = [x / norm for x in vec]
        return vec

    def embed_documents(self, texts):
        return [self.embed_text(t) for t in texts]

class MockOCRService(BaseOCR):
    def __init__(self, mock_text: str = "EXTRACTED TEST OCR TEXT 12345"):
        self.mock_text = mock_text

    def is_available(self) -> bool:
        return True

    def ocr_image(self, image_path: str) -> str:
        return self.mock_text

    def ocr_pdf(self, file_path: str):
        return {
            "document": os.path.basename(file_path),
            "pages": [{"page_number": 1, "text": self.mock_text}]
        }

class TestUniversalIngestion(unittest.TestCase):
    def setUp(self):
        from backend.app.config.settings import settings
        self.temp_dir = tempfile.TemporaryDirectory()
        self.db_fd, self.test_db_path = tempfile.mkstemp(suffix=".db")
        os.close(self.db_fd)
        self._orig_auth_db_path = settings.AUTH_DB_PATH
        settings.AUTH_DB_PATH = self.test_db_path
        os.environ["AUTH_DB_PATH"] = self.test_db_path
        init_db()

        # Isolated RAG Service
        self.rag_service = AegisRagService(
            persist_directory=os.path.join(self.temp_dir.name, "chroma"),
            embedding_model=MockEmbeddingModel(),
            safe_directories=[self.temp_dir.name],
            ocr_service=MockOCRService()
        )

    def tearDown(self):
        from backend.app.config.settings import settings
        self.temp_dir.cleanup()
        if os.path.exists(self.test_db_path):
            try:
                os.remove(self.test_db_path)
            except Exception:
                pass
        settings.AUTH_DB_PATH = self._orig_auth_db_path
        os.environ.pop("AUTH_DB_PATH", None)

    # -------------------------------------------------------------
    # 1. FILE DETECTOR & MAGIC BYTE SIGNATURE TESTS
    # -------------------------------------------------------------
    def test_detector_identifies_pdf(self):
        pdf_bytes = b"%PDF-1.5\n%\xe2\xe3\xcf\xd3\n1 0 obj\n<<>>\nendobj\n"
        res = FileDetector.detect_from_bytes(pdf_bytes, "manual.pdf")
        self.assertTrue(res.is_valid)
        self.assertTrue(res.is_safe)
        self.assertEqual(res.file_type, "pdf")
        self.assertEqual(res.category, "document")

    def test_detector_identifies_images(self):
        # PNG
        png_bytes = b"\x89PNG\r\n\x1a\n\x00\x00\x00\rIHDR"
        res_png = FileDetector.detect_from_bytes(png_bytes, "schematic.png")
        self.assertTrue(res_png.is_valid)
        self.assertEqual(res_png.category, "image")
        self.assertEqual(res_png.file_type, "png")

        # JPEG
        jpg_bytes = b"\xff\xd8\xff\xe0\x00\x10JFIF"
        res_jpg = FileDetector.detect_from_bytes(jpg_bytes, "pump.jpg")
        self.assertTrue(res_jpg.is_valid)
        self.assertEqual(res_jpg.category, "image")
        self.assertEqual(res_jpg.file_type, "jpeg")

    def test_detector_blocks_dangerous_binaries(self):
        # Linux ELF
        elf_bytes = b"\x7fELF\x02\x01\x01\x00" + b"\x00" * 30
        res_elf = FileDetector.detect_from_bytes(elf_bytes, "malicious.pdf")
        self.assertFalse(res_elf.is_safe)
        self.assertFalse(res_elf.is_valid)
        self.assertEqual(res_elf.category, "blocked")

        # Windows PE
        pe_bytes = bytearray(b"MZ" + b"\x00" * 0x3a + b"\x80\x00\x00\x00" + b"\x00" * 0x40)
        pe_bytes[0x80:0x84] = b"PE\x00\x00"
        res_pe = FileDetector.detect_from_bytes(bytes(pe_bytes), "tool.exe")
        self.assertFalse(res_pe.is_safe)
        self.assertFalse(res_pe.is_valid)

        # macOS Mach-O
        macho_bytes = b"\xfe\xed\xfa\xcf\x00\x00\x00\x01" + b"\x00" * 20
        res_macho = FileDetector.detect_from_bytes(macho_bytes, "runner.bin")
        self.assertFalse(res_macho.is_safe)

    def test_detector_identifies_code_and_text(self):
        py_bytes = b"def calculate_total(a, b):\n    return a + b\n"
        res_py = FileDetector.detect_from_bytes(py_bytes, "math_utils.py")
        self.assertTrue(res_py.is_valid)
        self.assertEqual(res_py.category, "code")
        self.assertEqual(res_py.file_type, "python")

        sql_bytes = b"SELECT * FROM refinery_metrics WHERE pressure > 100;"
        res_sql = FileDetector.detect_from_bytes(sql_bytes, "query.sql")
        self.assertTrue(res_sql.is_valid)
        self.assertEqual(res_sql.category, "code")

        csv_bytes = b"Sensor,Temperature,Pressure\nS1,45.2,101.3\nS2,48.1,102.5\n"
        res_csv = FileDetector.detect_from_bytes(csv_bytes, "telemetry.csv")
        self.assertTrue(res_csv.is_valid)
        self.assertEqual(res_csv.category, "spreadsheet")
        self.assertEqual(res_csv.file_type, "csv")

    # -------------------------------------------------------------
    # 2. EXTRACTORS FUNCTIONAL TESTS
    # -------------------------------------------------------------
    def test_docx_extractor(self):
        doc_path = os.path.join(self.temp_dir.name, "spec.docx")
        doc = docx.Document()
        doc.add_heading("Technical Specification", level=1)
        doc.add_paragraph("This document outlines sovereign industrial AI architectures.")
        table = doc.add_table(rows=2, cols=2)
        table.cell(0, 0).text = "Metric"
        table.cell(0, 1).text = "Target"
        table.cell(1, 0).text = "Latency"
        table.cell(1, 1).text = "<100ms"
        doc.save(doc_path)

        norm_doc = UniversalExtractorRegistry.extract_document(doc_path)
        self.assertEqual(norm_doc.category, "document")
        self.assertEqual(norm_doc.file_type, "docx")
        self.assertTrue(len(norm_doc.pages) >= 1)
        total_text = norm_doc.get_total_text()
        self.assertIn("Technical Specification", total_text)
        self.assertIn("Metric | Target", total_text)
        self.assertIn("Latency | <100ms", total_text)

    def test_excel_extractor(self):
        xlsx_path = os.path.join(self.temp_dir.name, "equipment.xlsx")
        wb = openpyxl.Workbook()
        ws1 = wb.active
        ws1.title = "Pumps"
        ws1.append(["Pump_ID", "Status", "Pressure_PSI"])
        ws1.append(["P-101", "OPERATIONAL", 145.2])
        ws1.append(["P-102", "MAINTENANCE", 0.0])

        ws2 = wb.create_sheet(title="Valves")
        ws2.append(["Valve_ID", "Position", "Flow_Rate"])
        ws2.append(["V-201", "OPEN", 85.5])
        wb.save(xlsx_path)

        norm_doc = UniversalExtractorRegistry.extract_document(xlsx_path)
        self.assertEqual(norm_doc.category, "spreadsheet")
        self.assertEqual(norm_doc.file_type, "xlsx")
        self.assertEqual(len(norm_doc.pages), 2)
        self.assertIn("Pumps", norm_doc.pages[0].section_title)
        self.assertIn("Valves", norm_doc.pages[1].section_title)
        self.assertIn("P-101", norm_doc.pages[0].text)
        self.assertIn("V-201", norm_doc.pages[1].text)

    def test_pptx_extractor(self):
        pptx_path = os.path.join(self.temp_dir.name, "presentation.pptx")
        prs = pptx.Presentation()
        slide_layout = prs.slide_layouts[0]
        slide1 = prs.slides.add_slide(slide_layout)
        title1 = slide1.shapes.title
        title1.text = "AEGIS Architecture Overview"
        subtitle = slide1.placeholders[1]
        subtitle.text = "Air-Gapped Sovereign AI Platform"

        slide2_layout = prs.slide_layouts[1]
        slide2 = prs.slides.add_slide(slide2_layout)
        slide2.shapes.title.text = "Key Invariants"
        slide2.placeholders[1].text = "100% On-Premise Execution\nDeterministic Audit HMAC Ledger"
        prs.save(pptx_path)

        norm_doc = UniversalExtractorRegistry.extract_document(pptx_path)
        self.assertEqual(norm_doc.category, "presentation")
        self.assertEqual(norm_doc.file_type, "pptx")
        self.assertEqual(len(norm_doc.pages), 2)
        self.assertIn("Slide 1", norm_doc.pages[0].section_title)
        self.assertIn("AEGIS Architecture Overview", norm_doc.pages[0].text)
        self.assertIn("Slide 2", norm_doc.pages[1].section_title)
        self.assertIn("Key Invariants", norm_doc.pages[1].text)

    def test_image_multimodal_extractor(self):
        img_path = os.path.join(self.temp_dir.name, "test_schematic.png")
        img = Image.new("RGB", (300, 200), color=(255, 255, 255))
        d = ImageDraw.Draw(img)
        d.text((10, 10), "PRESSURE SENSOR S-42", fill=(0, 0, 0))
        img.save(img_path)

        norm_doc = UniversalExtractorRegistry.extract_document(
            img_path,
            ocr_service=MockOCRService(mock_text="PRESSURE SENSOR S-42")
        )
        self.assertEqual(norm_doc.category, "image")
        self.assertEqual(norm_doc.file_type, "png")
        self.assertEqual(norm_doc.pages[0].metadata["width"], 300)
        self.assertEqual(norm_doc.pages[0].metadata["height"], 200)
        self.assertIn("PRESSURE SENSOR S-42", norm_doc.pages[0].text)

    def test_code_extractor(self):
        py_path = os.path.join(self.temp_dir.name, "pipeline.py")
        with open(py_path, "w", encoding="utf-8") as f:
            f.write("# AEGIS Controller\ndef run_pipeline(data):\n    return [x * 2 for x in data]\n")

        norm_doc = UniversalExtractorRegistry.extract_document(py_path)
        self.assertEqual(norm_doc.category, "code")
        self.assertEqual(norm_doc.file_type, "python")
        self.assertIn("def run_pipeline", norm_doc.pages[0].text)
        self.assertIn("Lines 1-", norm_doc.pages[0].section_title)

    # -------------------------------------------------------------
    # 3. END-TO-END RAG INGESTION & CITATIONS
    # -------------------------------------------------------------
    def test_e2e_universal_ingestion_and_search(self):
        # 1. Ingest PPTX
        pptx_path = os.path.join(self.temp_dir.name, "refinery_plan.pptx")
        prs = pptx.Presentation()
        s1 = prs.slides.add_slide(prs.slide_layouts[0])
        s1.shapes.title.text = "Turbine Maintenance Protocol"
        s1.placeholders[1].text = "Turbine T-900 must undergo vibration inspection every 48 operating hours."
        prs.save(pptx_path)

        doc_id_pptx = self.rag_service.ingest_document(pptx_path, owner_id=1, owner_username="engineer")
        self.assertTrue(len(doc_id_pptx) > 0)

        # 2. Ingest XLSX
        xlsx_path = os.path.join(self.temp_dir.name, "refinery_telemetry.xlsx")
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = "SensorReadings"
        ws.append(["Sensor_Tag", "Pressure_Bar", "Temperature_C", "Status"])
        ws.append(["PT-401", 12.4, 210.5, "NORMAL"])
        ws.append(["PT-402", 28.9, 345.1, "ALERT_OVERPRESSURE"])
        wb.save(xlsx_path)

        doc_id_xlsx = self.rag_service.ingest_document(xlsx_path, owner_id=1, owner_username="engineer")
        self.assertTrue(len(doc_id_xlsx) > 0)

        # 3. Ingest Python Code
        code_path = os.path.join(self.temp_dir.name, "safety_interlock.py")
        with open(code_path, "w", encoding="utf-8") as f:
            f.write("def check_emergency_shutdown(pressure_bar):\n    if pressure_bar > 25.0:\n        return True\n    return False\n")

        doc_id_code = self.rag_service.ingest_document(code_path, owner_id=1, owner_username="engineer")
        self.assertTrue(len(doc_id_code) > 0)

        # Search across ingested multimodal repository
        results_turbine = self.rag_service.search("Turbine vibration inspection", top_k=2)
        self.assertTrue(len(results_turbine) >= 1)
        self.assertEqual(results_turbine[0]["metadata"]["filename"], "refinery_plan.pptx")
        self.assertEqual(results_turbine[0]["metadata"]["category"], "presentation")
        self.assertIn("Slide 1", results_turbine[0]["metadata"]["section"])

        results_pressure = self.rag_service.search("ALERT_OVERPRESSURE PT-402", top_k=2)
        self.assertTrue(len(results_pressure) >= 1)
        self.assertEqual(results_pressure[0]["metadata"]["filename"], "refinery_telemetry.xlsx")
        self.assertEqual(results_pressure[0]["metadata"]["category"], "spreadsheet")
        self.assertIn("SensorReadings", results_pressure[0]["metadata"]["section"])

        # Test listing & stats
        docs = self.rag_service.list_documents(owner_id=1)
        self.assertEqual(len(docs), 3)
        categories = {d["category"] for d in docs}
        self.assertIn("presentation", categories)
        self.assertIn("spreadsheet", categories)
        self.assertIn("code", categories)

    def test_scanned_pdf_ocr_fallback(self):
        # Create a digital PDF with empty/scanned content (<15 characters)
        pdf_path = os.path.join(self.temp_dir.name, "scanned_invoice.pdf")
        import fitz
        doc = fitz.open()
        page = doc.new_page()
        # Page has no text
        doc.save(pdf_path)
        doc.close()

        norm_doc = UniversalExtractorRegistry.extract_document(
            pdf_path,
            ocr_service=MockOCRService(mock_text="INVOICE #98765 TOTAL: $4,500.00")
        )
        self.assertEqual(norm_doc.category, "document")
        self.assertEqual(norm_doc.file_type, "pdf")
        self.assertEqual(norm_doc.extraction_method, "ocr_scanned_pdf")
        self.assertTrue(len(norm_doc.pages) >= 1)
        self.assertIn("INVOICE #98765", norm_doc.pages[0].text)

    def test_markdown_and_json_ingestion(self):
        # Markdown
        md_path = os.path.join(self.temp_dir.name, "sop.md")
        with open(md_path, "w", encoding="utf-8") as f:
            f.write("# Standard Operating Procedure\n## Emergency Shutoff\nTurn valve V-1 clockwise.")

        norm_md = UniversalExtractorRegistry.extract_document(md_path)
        self.assertEqual(norm_md.category, "text")
        self.assertEqual(norm_md.file_type, "markdown")
        self.assertIn("Emergency Shutoff", norm_md.get_total_text())

        # JSON
        json_path = os.path.join(self.temp_dir.name, "config.json")
        with open(json_path, "w", encoding="utf-8") as f:
            json.dump({"refinery_id": "MRPL-01", "active_sensors": 42}, f)

        norm_json = UniversalExtractorRegistry.extract_document(json_path)
        self.assertEqual(norm_json.category, "code")
        self.assertEqual(norm_json.file_type, "json")
        self.assertIn("MRPL-01", norm_json.get_total_text())

    def test_mismatched_extension_safety(self):
        # File has .txt extension but contains ELF executable header
        fake_txt = os.path.join(self.temp_dir.name, "innocent.txt")
        with open(fake_txt, "wb") as f:
            f.write(b"\x7fELF\x02\x01\x01\x00" + b"\x00" * 50)

        with self.assertRaises(ValueError) as ctx:
            UniversalExtractorRegistry.extract_document(fake_txt)
        self.assertIn("Dangerous executable", str(ctx.exception))

    def test_corrupted_file_handling(self):
        corrupted_path = os.path.join(self.temp_dir.name, "broken.docx")
        with open(corrupted_path, "wb") as f:
            f.write(b"PK\x03\x04NOT_A_VALID_ZIP_ARCHIVE_DATA")

        with self.assertRaises(Exception):
            UniversalExtractorRegistry.extract_document(corrupted_path)

    def test_multi_user_document_scoping(self):
        # User 10 (Operator A)
        f1 = os.path.join(self.temp_dir.name, "operator_a_doc.txt")
        with open(f1, "w", encoding="utf-8") as f:
            f.write("Secret project alpha protocol for Operator A.")
        self.rag_service.ingest_document(f1, owner_id=10, owner_username="operator_a")

        # User 20 (Operator B)
        f2 = os.path.join(self.temp_dir.name, "operator_b_doc.txt")
        with open(f2, "w", encoding="utf-8") as f:
            f.write("Secret project beta protocol for Operator B.")
        self.rag_service.ingest_document(f2, owner_id=20, owner_username="operator_b")

        # Operator A listing
        docs_a = self.rag_service.list_documents(owner_id=10)
        self.assertEqual(len(docs_a), 1)
        self.assertEqual(docs_a[0]["filename"], "operator_a_doc.txt")

        # Operator B listing
        docs_b = self.rag_service.list_documents(owner_id=20)
        self.assertEqual(len(docs_b), 1)
        self.assertEqual(docs_b[0]["filename"], "operator_b_doc.txt")

        # Admin listing sees both
        docs_admin = self.rag_service.list_documents(is_admin=True)
        self.assertEqual(len(docs_admin), 2)

if __name__ == "__main__":
    unittest.main()
